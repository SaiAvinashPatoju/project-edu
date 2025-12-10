"""
Export service for generating PDF and PPTX presentations from slides
"""
import os
import json
import uuid
import tempfile
from typing import List, Dict, Any
from datetime import datetime
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

from models import Slide, LectureSession


class ExportService:
    """Service for exporting slides to PDF and PPTX formats"""
    
    def __init__(self):
        self.temp_dir = os.getenv("TEMP_FILE_DIR", tempfile.gettempdir())
        os.makedirs(self.temp_dir, exist_ok=True)
    
    def generate_pdf(self, slides: List[Slide], session: LectureSession) -> str:
        """
        Generate a PDF document from slides
        
        Args:
            slides: List of slide objects
            session: Lecture session object
            
        Returns:
            str: Path to the generated PDF file
        """
        # Create unique filename
        filename = f"lecture_slides_{session.id}_{uuid.uuid4().hex[:8]}.pdf"
        filepath = os.path.join(self.temp_dir, filename)
        
        # Create PDF document
        doc = SimpleDocTemplate(filepath, pagesize=letter)
        styles = getSampleStyleSheet()
        
        # Create custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=18,
            spaceAfter=30,
            alignment=1  # Center alignment
        )
        
        slide_title_style = ParagraphStyle(
            'SlideTitle',
            parent=styles['Heading2'],
            fontSize=16,
            spaceAfter=20,
            textColor='darkblue'
        )
        
        content_style = ParagraphStyle(
            'SlideContent',
            parent=styles['Normal'],
            fontSize=12,
            leftIndent=20,
            spaceAfter=10
        )
        
        # Build document content
        story = []
        
        # Add title page
        story.append(Paragraph(session.title or "Lecture Slides", title_style))
        story.append(Spacer(1, 0.5*inch))
        story.append(Paragraph(f"Generated on: {datetime.now().strftime('%B %d, %Y')}", styles['Normal']))
        story.append(Spacer(1, 0.5*inch))
        
        # Add slides
        for slide in slides:
            # Add slide title
            story.append(Paragraph(f"Slide {slide.slide_number}: {slide.title}", slide_title_style))
            
            # Parse and add content
            try:
                content_items = json.loads(slide.content) if isinstance(slide.content, str) else slide.content
                if isinstance(content_items, list):
                    for item in content_items:
                        story.append(Paragraph(f"• {item}", content_style))
                else:
                    story.append(Paragraph(str(content_items), content_style))
            except (json.JSONDecodeError, TypeError):
                # Fallback for non-JSON content
                story.append(Paragraph(str(slide.content), content_style))
            
            story.append(Spacer(1, 0.3*inch))
        
        # Build PDF
        doc.build(story)
        
        return filepath
    
    def generate_pptx(self, slides: List[Slide], session: LectureSession) -> str:
        """
        Generate a PPTX presentation from slides
        
        Args:
            slides: List of slide objects
            session: Lecture session object
            
        Returns:
            str: Path to the generated PPTX file
        """
        # Create unique filename
        filename = f"lecture_slides_{session.id}_{uuid.uuid4().hex[:8]}.pptx"
        filepath = os.path.join(self.temp_dir, filename)
        
        # Create presentation
        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        title_slide = prs.slides.add_slide(title_slide_layout)
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        title.text = session.title or "Lecture Slides"
        subtitle.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
        
        # Add content slides
        for slide in slides:
            # Use title and content layout
            slide_layout = prs.slide_layouts[1]  # Title and content layout
            new_slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            title_shape = new_slide.shapes.title
            title_shape.text = slide.title
            
            # Set content
            content_shape = new_slide.placeholders[1]  # Content placeholder
            text_frame = content_shape.text_frame
            text_frame.clear()  # Clear default text
            
            try:
                content_items = json.loads(slide.content) if isinstance(slide.content, str) else slide.content
                if isinstance(content_items, list):
                    for i, item in enumerate(content_items):
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        p.text = str(item)
                        p.level = 0  # Bullet level
                else:
                    text_frame.paragraphs[0].text = str(content_items)
            except (json.JSONDecodeError, TypeError):
                # Fallback for non-JSON content
                text_frame.paragraphs[0].text = str(slide.content)
        
        # Save presentation
        prs.save(filepath)
        
        return filepath
    
    def cleanup_file(self, filepath: str) -> bool:
        """
        Clean up a generated file
        
        Args:
            filepath: Path to the file to delete
            
        Returns:
            bool: True if successful, False otherwise
        """
        try:
            if os.path.exists(filepath):
                os.remove(filepath)
                return True
            return False
        except Exception:
            return False


# Global instance
export_service = ExportService()